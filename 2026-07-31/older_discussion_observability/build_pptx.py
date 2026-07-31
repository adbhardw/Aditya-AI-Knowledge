from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── colour palette ──────────────────────────────────────────────
BG       = RGBColor(0x0D, 0x11, 0x17)   # near-black bg
BG2      = RGBColor(0x16, 0x1B, 0x22)   # card bg
BORDER   = RGBColor(0x30, 0x36, 0x3D)   # subtle border
WHITE    = RGBColor(0xE6, 0xED, 0xF3)
GRAY     = RGBColor(0x8B, 0x94, 0x9E)
BLUE     = RGBColor(0x58, 0xA6, 0xFF)   # V1
ORANGE   = RGBColor(0xF0, 0x88, 0x3E)   # V2
GREEN    = RGBColor(0x3F, 0xB9, 0x50)   # V3
RED      = RGBColor(0xFF, 0x7B, 0x72)   # failure
PURPLE   = RGBColor(0xD2, 0xA8, 0xFF)   # arch
DARKBLUE = RGBColor(0x1F, 0x6F, 0xEB)
DARKORNG = RGBColor(0xE3, 0x67, 0x23)
DARKGRN  = RGBColor(0x2E, 0xA0, 0x43)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]   # completely blank

# ── helpers ─────────────────────────────────────────────────────

def add_slide():
    s = prs.slides.add_slide(blank_layout)
    fill_bg(s, BG)
    return s

def fill_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def box(slide, x, y, w, h, bg=BG2, border=BORDER, radius=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    shape.line.color.rgb = border
    shape.line.width = Pt(0.75)
    return shape

def left_bar(slide, x, y, h, color):
    shape = slide.shapes.add_shape(
        1, Inches(x), Inches(y), Inches(0.07), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txt(slide, text, x, y, w, h,
        size=12, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb

def heading(slide, title, subtitle=None, chip=None, chip_color=BLUE):
    # colour accent bar at top
    bar = slide.shapes.add_shape(1, 0, 0, W, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = chip_color
    bar.line.fill.background()

    if chip:
        box(slide, 0.45, 0.18, 1.6, 0.26, bg=BG2, border=chip_color)
        txt(slide, chip, 0.45, 0.18, 1.6, 0.26,
            size=9, bold=True, color=chip_color, align=PP_ALIGN.CENTER)

    ytitle = 0.52 if chip else 0.22
    txt(slide, title, 0.45, ytitle, 12.4, 0.7,
        size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.45, ytitle+0.65, 12.4, 0.4,
            size=13, color=GRAY)

def sprint_bar(slide):
    y = 6.95
    shapes = [
        (DARKBLUE, "V1  ·  Sprints 1–2  ·  Activity Feed", 0.45, 3.8),
        (DARKORNG, "V2  ·  Sprints 3–5  ·  Webhook Delivery", 4.28, 4.9),
        (DARKGRN,  "V3  ·  Sprints 6–7  ·  Hardening",   9.21, 3.67),
    ]
    for color, label, x, w in shapes:
        b = slide.shapes.add_shape(1, Inches(x), Inches(y),
                                   Inches(w), Inches(0.38))
        b.fill.solid()
        b.fill.fore_color.rgb = color
        b.line.fill.background()
        txt(slide, label, x, y, w, 0.38,
            size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

def card(slide, x, y, w, h, title, body_lines,
         accent=BLUE, tag=None):
    box(slide, x, y, w, h)
    left_bar(slide, x, y+0.05, h-0.1, accent)
    ty = y + 0.12
    if tag:
        txt(slide, tag, x+0.18, ty, w-0.25, 0.22,
            size=8.5, bold=True, color=accent)
        ty += 0.22
    txt(slide, title, x+0.18, ty, w-0.25, 0.28,
        size=11, bold=True, color=WHITE)
    ty += 0.28
    body = "\n".join(body_lines)
    txt(slide, body, x+0.18, ty, w-0.25, h-ty+y-0.1,
        size=9.5, color=GRAY)

def mono(slide, text, x, y, w, h):
    b = box(slide, x, y, w, h, bg=BG)
    tf_box = slide.shapes.add_textbox(
        Inches(x+0.12), Inches(y+0.1),
        Inches(w-0.2),  Inches(h-0.15))
    tf = tf_box.text_frame
    tf.word_wrap = False
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(9)
        run.font.color.rgb = RGBColor(0xE6,0xED,0xF3)
        try:
            run.font.name = "Courier New"
        except Exception:
            pass
    return b

# ════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════
s = add_slide()

# gradient-like accent block
accent = s.shapes.add_shape(1, 0, Inches(5.5), W, Inches(0.06))
accent.fill.solid(); accent.fill.fore_color.rgb = DARKBLUE
accent.line.fill.background()

txt(s, "DV-13856  ·  CFI-1388", 0.55, 0.9, 12, 0.4,
    size=11, color=GRAY)
txt(s, "Clean Room Observability\n& Callback Registration", 0.55, 1.3, 12, 1.9,
    size=40, bold=True, color=WHITE)
txt(s,
    "Platform-grade event delivery backbone for Habu.\n"
    "V1 ships an Activity Feed to every user.  V2 adds push webhooks to XMI and partners.\n"
    "V3 hardens the system for enterprise production.",
    0.55, 3.35, 11, 1.1, size=14, color=GRAY)

# sprint bar with labels below
for color, label, x, w in [
    (DARKBLUE, "V1 · Sprints 1–2",  0.55, 3.8),
    (DARKORNG, "V2 · Sprints 3–5",  4.38, 4.9),
    (DARKGRN,  "V3 · Sprints 6–7",  9.31, 3.47),
]:
    b = s.shapes.add_shape(1, Inches(x), Inches(5.7), Inches(w), Inches(0.45))
    b.fill.solid(); b.fill.fore_color.rgb = color
    b.line.fill.background()
    txt(s, label, x, 5.7, w, 0.45,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "Author: Aditya Bhardwaj  ·  May 2026  ·  For TLM + Product Review",
    0.55, 6.9, 12, 0.35, size=10, color=GRAY)

# ════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM
# ════════════════════════════════════════════════════════════════
s = add_slide()
heading(s, "The Problem — Objects Change, Nobody Knows",
        "External systems (XMI UI, customer integrations) hold stale references. "
        "No notification mechanism exists today.",
        chip="PROBLEM", chip_color=RED)

card(s, 0.45, 1.35, 3.9, 2.0, "unhygienix — Questions",
     ["Columns removed from output → dashboards break silently",
      "Dataset assignment swapped → wrong data imported",
      "Question deleted → 404s in downstream workflows",
      "Status DRAFT→PUBLISHED → integrations miss the event"],
     accent=BLUE, tag="QUESTION MUTATIONS")

card(s, 4.58, 1.35, 3.9, 2.0, "forebitt — Data Connections",
     ["Field type changed: Integer→String → SUM() fails at runtime",
      "Field added or removed → schema expectations violated",
      "Data connection deleted → every question using it breaks",
      "No warning sent to consumers before change takes effect"],
     accent=ORANGE, tag="DATA CONNECTION MUTATIONS")

card(s, 8.71, 1.35, 3.9, 2.0, "picanmix — Exports",
     ["Export job completion invisible to integrations",
      "Config changes undetected until next poll",
      "External systems poll every 30s — 99% return no change",
      "Reactive posture — not enterprise software behaviour"],
     accent=GREEN, tag="EXPORT MUTATIONS")

# bottom context box
b = box(s, 0.45, 3.55, 12.4, 0.95, bg=RGBColor(0x1F,0x6F,0xEB,),
         border=DARKBLUE)
b.fill.fore_color.rgb = RGBColor(0x0D,0x1A,0x35)
txt(s, "Two audiences, two solutions:",
    0.65, 3.65, 5, 0.3, size=11, bold=True, color=WHITE)
txt(s, "Human users on Habu  →  Activity Feed UI showing what changed, when, who did it  (V1)",
    0.65, 3.95, 11.8, 0.3, size=10.5, color=BLUE)
txt(s, "Programmatic systems (XMI, SafeHaven)  →  Push webhook — Habu calls their endpoint automatically  (V2)",
    0.65, 4.25, 11.8, 0.3, size=10.5, color=ORANGE)

sprint_bar(s)

# ════════════════════════════════════════════════════════════════
# SLIDE 3 — V1 ACTIVITY FEED
# ════════════════════════════════════════════════════════════════
s = add_slide()
heading(s, "V1 — Activity Feed UI",
        "One new table. One new API. Ships to every Habu user with zero setup. "
        "Interceptors built here become the foundation for V2 webhooks.",
        chip="V1 · SPRINTS 1–2", chip_color=BLUE)

# left: architecture
card(s, 0.45, 1.35, 4.4, 1.25,
     "unhygienix — CRUD Interceptor",
     ["Snapshot state BEFORE mutation",
      "Execute mutation (delete / update)",
      "Diff before vs after → build ObjectEvent",
      "[IN SAME DB TRANSACTION] INSERT object_events"],
     accent=BLUE, tag="NEW CODE — unhygienix")

txt(s, "↓  atomic  (rolls back if mutation rolls back)", 1.8, 2.65, 3.5, 0.3,
    size=9, color=GRAY)

card(s, 0.45, 3.0, 4.4, 1.3,
     "object_events table  (PostgreSQL)",
     ["event_id · cleanroom_id · object_type · object_id",
      "object_name · event_type · changed_fields JSONB",
      "performed_by · org_id · event_time",
      "Append-only — never updated, only inserted"],
     accent=BLUE, tag="NEW TABLE — unhygienix DB")

txt(s, "↓  GET /cleanrooms/{id}/events  →  external-api-server",
    1.5, 4.35, 4.0, 0.3, size=9, color=GRAY)

# right: UI mockup
box(s, 5.1, 1.35, 7.75, 5.1, bg=BG2, border=BORDER)
txt(s, "Activity Feed Tab — Clean Room UI", 5.2, 1.42, 7.5, 0.3,
    size=10, bold=True, color=GRAY)
box(s, 5.2, 1.73, 7.55, 0.3, bg=RGBColor(0x21,0x26,0x2D), border=BORDER)
txt(s, "Filter:  [All Types ▾]   [All Users ▾]   [Last 30 days ▾]",
    5.3, 1.75, 7.3, 0.25, size=9, color=GRAY)

events = [
    (BLUE,   "QUESTION UPDATED",       '"Revenue Analysis" — columns removed',
             "Removed: customer_segment, channel  ·  sarah@acme.com  ·  2026-05-12 14:32"),
    (ORANGE, "DATA CONNECTION UPDATED", '"CRM_Import" — field type changed',
             "purchase_count: Integer → String  ·  john@acme.com  ·  2026-05-12 11:15"),
    (RED,    "QUESTION DELETED",        '"Lookalike Segments" — permanently deleted',
             "admin@liveramp.com  ·  2026-05-11 16:45"),
    (GREEN,  "QUESTION UPDATED",        '"Attribution Model" — dataset swapped',
             "Q1_Sales_2026  →  Q2_Sales_2026  ·  sarah@acme.com  ·  2026-05-11 09:00"),
]
ey = 2.1
for color, etype, title, detail in events:
    left_bar(s, 5.2, ey, 0.72, color)
    txt(s, etype,  5.35, ey+0.03, 3.5, 0.22, size=8.5, bold=True, color=color)
    txt(s, title,  5.35, ey+0.24, 7.3, 0.24, size=10, bold=False, color=WHITE)
    txt(s, detail, 5.35, ey+0.47, 7.3, 0.2,  size=8.5, color=GRAY)
    ey += 0.82

sprint_bar(s)

# ════════════════════════════════════════════════════════════════
# SLIDE 4 — V2 WEBHOOK FLOW
# ════════════════════════════════════════════════════════════════
s = add_slide()
heading(s, "V2 — Webhook Callback Delivery",
        "unhygienix / forebitt / picanmix publish events. "
        "Pegleg worker fans out to XMI, SafeHaven, and future consumers independently.",
        chip="V2 · SPRINTS 3–5", chip_color=ORANGE)

# flow column
flow_items = [
    (ORANGE, "unhygienix  /  forebitt  /  picanmix",
     "Mutation fires → interceptor runs → INSERT object_events [in txn]\n[POST COMMIT async ~1ms] → SNS publish to habu-object-events topic"),
    (ORANGE, "SNS Topic: habu-object-events  (picanmix infra)",
     "Fan-out: each subscriber gets an independent copy\nAdding new consumer = new SQS subscription, zero code change"),
    (ORANGE, "SQS-XMI  ·  SQS-SafeHaven  ·  SQS-Future  (independent queues)",
     "Each queue is isolated — XMI being slow/down does not affect SafeHaven\nmaxReceiveCount=3 → auto-move to DLQ on failure"),
    (ORANGE, "Pegleg — Callback Delivery Worker  (Java, horizontally scalable)",
     "receiveMessage → query callback_registrations → build payload\nHMAC-SHA256 sign → circuit breaker check → HTTP POST\nwrite delivery_log → deleteMessage (ACK SQS)"),
    (GREEN,  "XMI  /  SafeHaven  /  Future Customer",
     "Verifies X-Habu-Signature · Updates UI · Refreshes cache · Triggers workflow\nReturns 200 OK → worker ACKs SQS"),
]
fy = 1.35
for color, title, body in flow_items:
    card(s, 0.45, fy, 6.5, 1.05, title, body.split("\n"),
         accent=color)
    if fy < 5.3:
        txt(s, "↓", 3.35, fy+1.06, 0.5, 0.22, size=14, color=GRAY,
            align=PP_ALIGN.CENTER)
    fy += 1.15

# right: registration + payload
card(s, 7.2, 1.35, 5.9, 2.1,
     "Registration — one-time setup by XMI",
     ["POST /cleanrooms/{crId}/callbacks",
      'objectType: "QUESTION"  ·  objectId: null (all questions)',
      "callbackUrl: https://xmi.liveramp.com/habu-events",
      "authConfig: { authType: BEARER }",
      "Response includes signingSecret (whsec_...) for HMAC verification"],
     accent=ORANGE, tag="NEW API — external-api-server")

card(s, 7.2, 3.65, 5.9, 2.55,
     "Payload XMI receives on every change",
     ['POST  /habu-events',
      'X-Habu-Signature: sha256=a3f9...  (V3)',
      'X-Habu-Event-Id: evt-789  (idempotency)',
      '---',
      'eventType: UPDATED  ·  objectType: QUESTION',
      'objectName: "Revenue Analysis"',
      'changedFields: { dimensions: { removed: [...] } }',
      'changedBy: sarah@acme.com'],
     accent=ORANGE, tag="WEBHOOK PAYLOAD")

sprint_bar(s)

# ════════════════════════════════════════════════════════════════
# SLIDE 5 — V3 HARDENING
# ════════════════════════════════════════════════════════════════
s = add_slide()
heading(s, "V3 — Production Hardening",
        "No new tables. Extends Pegleg delivery worker + adds Webhooks UI tab. "
        "Makes the system cryptographically secure, self-observable, and replayable.",
        chip="V3 · SPRINTS 6–7", chip_color=GREEN)

cards_v3 = [
    ("HMAC Signature\nVerification",
     ["X-Habu-Signature: sha256=... on every POST",
      "Per-registration secret — XMI ≠ SafeHaven",
      "X-Habu-Timestamp header prevents replay attacks",
      "XMI verifies on receipt, rejects if mismatch"],
     GREEN),
    ("Circuit Breaker\n(Resilience4j)",
     ["CLOSED → normal, calling endpoint",
      "OPEN → 5+ failures, stop calling 60s, free threads",
      "HALF-OPEN → 1 probe call, recover if 200 OK",
      "Shared state across workers via Redis"],
     GREEN),
    ("Dead Letter Queue\n+ Admin Replay UI",
     ["SQS auto-moves to DLQ after 3 failed attempts",
      "CloudWatch alarm: DLQ depth > 0 → PagerDuty",
      "New Webhooks tab: inspect failed messages",
      "One-click Replay button per failed event"],
     GREEN),
    ("Delivery\nObservability",
     ["delivery_log: every attempt, status, latency",
      "CloudWatch: success/failure count, latency histogram",
      "Per-registration: success rate, last delivered",
      "Circuit open count by endpoint URL"],
     GREEN),
    ("Webhook Test\nEndpoint",
     ["POST /callbacks/{id}/test",
      "Sends synthetic event to registered URL",
      "Returns http_status, response_ms, success",
      "XMI validates before going live"],
     GREEN),
    ("Event Replay\n& Idempotency",
     ["Replay API: re-deliver from object_events table",
      "X-Habu-Event-Id header on every POST",
      "Worker checks delivery_log before POSTing",
      "SQS at-least-once + two-layer dedup"],
     GREEN),
]

positions = [
    (0.45, 1.35), (4.58, 1.35), (8.71, 1.35),
    (0.45, 3.85), (4.58, 3.85), (8.71, 3.85),
]
for (cx, cy), (title, body, color) in zip(positions, cards_v3):
    card(s, cx, cy, 3.9, 2.25, title, body, accent=color)

sprint_bar(s)

# ════════════════════════════════════════════════════════════════
# SLIDE 6 — DATA MODEL & APIs
# ════════════════════════════════════════════════════════════════
s = add_slide()
heading(s, "Data Model — 3 New Tables + API Inventory",
        "V1 adds 1 table. V2 adds 2 tables. V3 adds none. All via Flyway migrations in unhygienix DB.",
        chip="DATA MODEL", chip_color=PURPLE)

# table headers
hdr_y = 1.35
col_x = [0.45, 2.85, 4.5,  7.3,  10.0]
col_w = [2.3,  1.55, 2.7,  2.6,   2.6]
col_h = [("Table", WHITE), ("Version", WHITE),
         ("Purpose", WHITE), ("Key Fields", WHITE), ("Written by", WHITE)]

for (label, color), x, w in zip(col_h, col_x, col_w):
    b = box(s, x, hdr_y, w, 0.32, bg=RGBColor(0x21,0x26,0x2D), border=BORDER)
    txt(s, label, x+0.08, hdr_y+0.05, w-0.1, 0.22,
        size=9.5, bold=True, color=GRAY)

rows = [
    ("object_events",          "V1", BLUE,
     "Atomic audit log of all mutations. Powers Activity Feed UI.",
     "event_id, object_type, event_type, changed_fields JSONB, performed_by",
     "unhygienix interceptor — inside transaction"),
    ("callback_registrations", "V2", ORANGE,
     "Stores registered callback URLs, auth, signing secrets per cleanroom.",
     "cleanroom_id, object_type, object_id, callback_url, signing_secret",
     "POST /callbacks API — external-api-server"),
    ("delivery_log",           "V2", ORANGE,
     "Records every HTTP delivery attempt. Powers idempotency + V3 UI.",
     "registration_id, event_id, attempt_number, http_status, response_ms",
     "Pegleg delivery worker — before ACK"),
]
ry = hdr_y + 0.34
for (tname, ver, vcolor, purpose, fields, writer) in rows:
    row_h = 0.75
    for x, w in zip(col_x, col_w):
        box(s, x, ry, w, row_h, bg=BG2, border=BORDER)
    left_bar(s, col_x[0], ry+0.05, row_h-0.1, vcolor)
    txt(s, tname,   col_x[0]+0.18, ry+0.08, col_w[0]-0.25, 0.58,
        size=9.5, bold=True, color=vcolor)
    txt(s, ver,     col_x[1]+0.08, ry+0.08, col_w[1]-0.12, 0.58,
        size=9.5, bold=True, color=vcolor)
    txt(s, purpose, col_x[2]+0.08, ry+0.08, col_w[2]-0.12, 0.65,
        size=9, color=GRAY)
    txt(s, fields,  col_x[3]+0.08, ry+0.08, col_w[3]-0.12, 0.65,
        size=8.5, color=GRAY)
    txt(s, writer,  col_x[4]+0.08, ry+0.08, col_w[4]-0.12, 0.65,
        size=8.5, color=GRAY)
    ry += row_h + 0.04

# API list
txt(s, "API INVENTORY", 0.45, ry+0.1, 12, 0.25,
    size=9.5, bold=True, color=PURPLE)
apis = [
    ("V1", "GET",    "/cleanrooms/{id}/events",               "external-api-server", "Activity feed read — since, objectType, performedBy filters"),
    ("V2", "POST",   "/cleanrooms/{id}/callbacks",            "external-api-server", "Register callback URL"),
    ("V2", "GET",    "/cleanrooms/{id}/callbacks",            "external-api-server", "List registrations"),
    ("V2", "DELETE", "/cleanrooms/{id}/callbacks/{cbId}",     "external-api-server", "Deregister"),
    ("V3", "POST",   "/cleanrooms/{id}/callbacks/{cbId}/test","external-api-server", "Send synthetic test event"),
    ("V3", "POST",   "/admin/callbacks/{cbId}/replay-dlq",    "external-api-server", "Replay DLQ messages"),
]
api_y = ry + 0.38
for ver, method, path, svc, desc in apis:
    vcolor = BLUE if ver=="V1" else ORANGE if ver=="V2" else GREEN
    mcolor = GREEN if method=="GET" else ORANGE if method=="POST" else RED
    txt(s, ver,    0.45, api_y, 0.4,  0.22, size=8.5, bold=True, color=vcolor)
    txt(s, method, 0.9,  api_y, 0.65, 0.22, size=8.5, bold=True, color=mcolor)
    txt(s, path,   1.6,  api_y, 4.2,  0.22, size=8.5, color=WHITE)
    txt(s, svc,    5.9,  api_y, 2.8,  0.22, size=8.5, color=GRAY)
    txt(s, desc,   8.8,  api_y, 4.3,  0.22, size=8.5, color=GRAY)
    api_y += 0.27

sprint_bar(s)

# ════════════════════════════════════════════════════════════════
# SLIDE 7 — COMBINED ARCHITECTURE
# ════════════════════════════════════════════════════════════════
s = add_slide()
heading(s, "Combined Architecture — V1 + V2 + V3",
        "Single interceptor. Atomic audit log. Fan-out delivery to N independent consumers. "
        "Pegleg worker handles all webhook delivery with retry, DLQ, and circuit breaker.",
        chip="FULL DESIGN", chip_color=PURPLE)

arch_text = """\
  unhygienix (Questions)    forebitt (Data Connections)    picanmix (Exports)
        |                           |                             |
        +---------------------------+-----------------------------+
                                    |
                    ┌───────────────▼───────────────┐
                    │    CRUD Interceptor Layer      │
                    │  snapshot → execute → diff     │
                    └───────┬───────────────┬────────┘
         [IN TRANSACTION]   │        [POST COMMIT async ~1ms]
                            ▼                   ▼
                   object_events table    SNS habu-object-events
                   (V1 audit log)                │
                        │               ┌────────┼─────────┐
                   UI reads here        ▼        ▼         ▼
                   GET /events     SQS-XMI  SQS-SafeHvn  SQS-Future
                   Activity Feed        │    (independent, isolated)
                                        ▼
                             ┌──────────────────────────┐
                             │  Pegleg Delivery Worker   │
                             │  (Java, N instances)      │
                             │  1. query callback_regs   │
                             │  2. build JSON payload    │
                             │  3. HMAC-SHA256 sign (V3) │
                             │  4. circuit breaker  (V3) │
                             │  5. HTTP POST             │
                             │  6. write delivery_log    │
                             │  7. ACK SQS               │
                             └─────────┬────────┬────────┘
                                    200 OK   failure → retry 1s→5s→30s
                                       │             │
                               XMI / SafeHaven      DLQ ──▶ CloudWatch
                               verifies HMAC        DLQ ──▶ PagerDuty
                               updates UI           DLQ ──▶ Admin Replay UI (V3)"""

mono(s, arch_text, 0.35, 1.35, 12.6, 5.45)
sprint_bar(s)

# ════════════════════════════════════════════════════════════════
# SLIDE 8 — SPRINT ROADMAP
# ════════════════════════════════════════════════════════════════
s = add_slide()
heading(s, "Sprint Roadmap",
        "V1 ships value to every user in 2 sprints with zero infra risk. "
        "V2 and V3 build on top without touching V1 code.",
        chip="ROADMAP", chip_color=BLUE)

roadmap = [
    ("V1 · SPRINTS 1–2 · Activity Feed UI", BLUE, DARKBLUE, [
        ("NEW TABLE",    "object_events — unhygienix DB — Flyway migration"),
        ("NEW CODE",     "CRUD interceptors + field differ — unhygienix"),
        ("NEW CODE",     "GET /cleanrooms/{id}/events — external-api-server"),
        ("NEW UI",       "Activity Feed tab with filters — Frontend"),
    ]),
    ("V2 · SPRINTS 3–5 · Webhook Callback Delivery", ORANGE, DARKORNG, [
        ("NEW TABLE",    "callback_registrations — unhygienix DB"),
        ("NEW TABLE",    "delivery_log — unhygienix DB"),
        ("NEW INFRA",    "SNS topic + SQS queues (per consumer) + DLQs — picanmix"),
        ("NEW INFRA",    "CloudWatch alarm on DLQ depth > 0"),
        ("EXTEND",       "SNS publish added to V1 interceptors — unhygienix"),
        ("NEW CODE",     "Callback CRUD APIs — external-api-server"),
        ("NEW SERVICE",  "Pegleg — Callback Delivery Worker (Java) — SQS consumer + POST + retry"),
        ("NEW PROTOS",   "CallbackRegistration · CallbackEvent · CallbackAuthConfig — ignoramus"),
    ]),
    ("V3 · SPRINTS 6–7 · Production Hardening", GREEN, DARKGRN, [
        ("EXTEND",       "HMAC signature + timestamp replay prevention — Pegleg worker"),
        ("EXTEND",       "Circuit breaker Resilience4j per endpoint — Pegleg worker"),
        ("EXTEND",       "Per-field subscription filter — Pegleg worker"),
        ("NEW API",      "POST /callbacks/{id}/test — webhook test endpoint"),
        ("NEW API",      "POST /admin/callbacks/{cbId}/replay-dlq"),
        ("NEW UI",       "Webhooks tab — delivery status + DLQ inspect + Replay — Frontend"),
    ]),
]

ry = 1.35
for phase_title, color, dark_color, items in roadmap:
    ph = 0.3 + len(items) * 0.32
    b = s.shapes.add_shape(1, Inches(0.45), Inches(ry),
                           Inches(12.43), Inches(ph))
    b.fill.solid(); b.fill.fore_color.rgb = dark_color
    b.line.color.rgb = color; b.line.width = Pt(0.75)

    txt(s, phase_title, 0.6, ry+0.04, 12, 0.25,
        size=11, bold=True, color=WHITE)

    iy = ry + 0.31
    for badge, desc in items:
        bc = BLUE if "TABLE" in badge else ORANGE if "INFRA" in badge or "PROTO" in badge \
             else GREEN if "EXTEND" in badge else PURPLE if "SERVICE" in badge or "API" in badge \
             else GRAY
        txt(s, f"[{badge}]", 0.65, iy, 1.45, 0.25, size=8, bold=True, color=bc)
        txt(s, desc, 2.15, iy, 10.5, 0.25, size=9.5, color=WHITE)
        iy += 0.30

    ry += ph + 0.1

sprint_bar(s)

# ── save ────────────────────────────────────────────────────────
out = "/Users/adbhar/Documents/project-sprint-understanding-1/aditya_observability_discussion/DV13856_Observability_Presentation.pptx"
prs.save(out)
print("Saved:", out)
